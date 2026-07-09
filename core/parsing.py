import os
import fitz  # PyMuPDF

def extract_text_from_file(file) -> str:
    """
    Extracts raw text from an uploaded file based on its extension.
    Supports PDF, MD, TXT, DOCX, PPTX, CSV, XLSX, HTML, and HTM.
    
    Args:
        file: A file-like object from st.file_uploader.
        
    Returns:
        str: The extracted raw text, structured appropriately.
        
    Raises:
        ValueError: If the file type is unsupported or parsing fails.
        Exception: For any errors during parsing.
    """
    filename = file.name
    _, ext = os.path.splitext(filename.lower())
    
    if ext in [".md", ".txt"]:
        try:
            content = file.read()
            if isinstance(content, bytes):
                return content.decode("utf-8")
            return content
        except Exception as e:
            raise Exception(f"Failed to read text file: {str(e)}")
            
    elif ext == ".pdf":
        try:
            pdf_bytes = file.read()
            doc = fitz.open(stream=pdf_bytes, filetype="pdf")
            text_parts = []
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                text_parts.append(page.get_text())
            doc.close()
            extracted_text = "\n".join(text_parts)
            if not extracted_text.strip():
                raise ValueError("PDF document contains no selectable text.")
            return extracted_text
        except Exception as e:
            raise Exception(f"Failed to parse PDF document: {str(e)}")
            
    elif ext == ".docx":
        try:
            import docx
            doc = docx.Document(file)
            lines = [p.text for p in doc.paragraphs if p.text.strip()]
            return "\n\n".join(lines)
        except Exception as e:
            raise Exception(f"Failed to parse Word Document (.docx): {str(e)}")
            
    elif ext == ".pptx":
        try:
            from pptx import Presentation
            prs = Presentation(file)
            slides_text = []
            for idx, slide in enumerate(prs.slides):
                slide_parts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_parts.append(shape.text.strip())
                if slide_parts:
                    slides_text.append(f"## Slide {idx + 1}\n\n" + "\n".join(slide_parts))
            return "\n\n".join(slides_text)
        except Exception as e:
            raise Exception(f"Failed to parse PowerPoint Presentation (.pptx): {str(e)}")
            
    elif ext == ".csv":
        try:
            import pandas as pd
            df = pd.read_csv(file)
            # Cap data row budget at 5,000 to safe-guard server memory
            truncated_df = df.head(5000)
            if truncated_df.empty:
                return "Empty CSV table."
            return truncated_df.to_markdown(index=False)
        except Exception as e:
            raise Exception(f"Failed to parse CSV file: {str(e)}")
            
    elif ext == ".xlsx":
        try:
            import pandas as pd
            df = pd.read_excel(file)
            # Cap data row budget at 5,000 to safe-guard server memory
            truncated_df = df.head(5000)
            if truncated_df.empty:
                return "Empty Excel table."
            return truncated_df.to_markdown(index=False)
        except Exception as e:
            raise Exception(f"Failed to parse Excel Spreadsheet (.xlsx): {str(e)}")

    elif ext in [".html", ".htm"]:
        try:
            from bs4 import BeautifulSoup
            content = file.read()
            if isinstance(content, bytes):
                html_content = content.decode("utf-8")
            else:
                html_content = content
            
            soup = BeautifulSoup(html_content, "html.parser")
            
            # Decompose boilerplate and interactive elements in place
            noise_elements = ['header', 'footer', 'nav', 'aside', 'script', 'style', 'form', 'button', 'iframe', 'noscript']
            for tag in noise_elements:
                for element in soup.find_all(tag):
                    element.decompose()
            
            # Extract clean text separator mapping
            extracted_text = soup.get_text(separator="\n")
            
            # Sanitize extreme padding and repetitive empty lines
            lines = [line.strip() for line in extracted_text.splitlines() if line.strip()]
            return "\n".join(lines)
        except Exception as e:
            raise Exception(f"Failed to parse HTML document: {str(e)}")
            
    else:
        raise ValueError(
            f"Unsupported file type '{ext}'. Supported extensions: .pdf, .md, .txt, .docx, .pptx, .csv, .xlsx, .html, .htm"
        )
